import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx not installed. Please run: pip install python-pptx")
    sys.exit(1)

def create_deck():
    prs = Presentation()
    
    # Define colors
    bg_color = RGBColor(18, 18, 36)      # midnight navy
    title_color = RGBColor(245, 200, 66)  # gold
    text_color = RGBColor(242, 242, 250)   # white-ish
    muted_color = RGBColor(154, 160, 199) # grey-blue
    
    # Helper to set slide background
    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    # Helper to add a formatted title
    def add_title(slide, text):
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = title_color

    # ----------------------------------------------------
    # Slide 1: Title & Problem Statement
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    
    # Title
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BIOME: IB ESS Exam Arena"
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = title_color
    
    p2 = tf.add_paragraph()
    p2.text = "Reimagining Written Assessment Preparation for Environmental Systems"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(18)
    p2.font.color.rgb = muted_color
    p2.space_before = Pt(10)

    # Problem Statement Section
    tx_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(4))
    tf2 = tx_box2.text_frame
    tf2.word_wrap = True
    
    p3 = tf2.paragraphs[0]
    p3.text = "⚠️ THE PROBLEM STATEMENT"
    p3.font.name = 'Trebuchet MS'
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = title_color
    p3.space_after = Pt(10)
    
    bullets = [
        "Traditional study companions rely heavily on multiple-choice questions (MCQs).",
        "However, the official IB DP ESS curriculum has ZERO multiple-choice questions in its written examinations.",
        "Students are assessed entirely on structured data analysis case studies (Paper 1) and evaluative 9-mark essays (Paper 2).",
        "Without proper platforms to draft, revise, and grade written answers against official rubrics, students remain unprepared for actual exam expectations."
    ]
    for bullet in bullets:
        p = tf2.add_paragraph()
        p.text = "• " + bullet
        p.font.name = 'Calibri'
        p.font.size = Pt(16)
        p.font.color.rgb = text_color
        p.space_after = Pt(12)
        
    # ----------------------------------------------------
    # Slide 2: The Solution - BIOME Exam Simulator
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_title(slide, "The Solution: BIOME Exam Simulator")
    
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    points = [
        ("Paper 1 Simulator", "Interactive data response modules based on authentic IB case studies (e.g., The Salton Sea System Collapse) featuring systems flow maps and structured multi-mark questions."),
        ("Paper 2 Section A Simulator", "Thermodynamics and feedback loop questions requiring scientific explanations of feedback cycles (e.g., Permafrost Methane positive loop)."),
        ("Paper 2 Section B Essay Arena", "Full-length 9-mark evaluative essay drafting zone integrating the Brundtland definition and official essay structure guidance.")
    ]
    for title, desc in points:
        p = tf.add_paragraph()
        p.text = f"✨ {title}"
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_before = Pt(8)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Calibri'
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = text_color
        p_desc.space_after = Pt(10)

    # ----------------------------------------------------
    # Slide 3: Interactive AI Tutor Grading
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_title(slide, "Real-Time AI Tutor Grading & Rubrics")
    
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    features = [
        ("Submit for Grading", "Students write structured answers inside textareas and submit them directly to the AI Tutor with a single click."),
        ("Official Markscheme Alignment", "The frontend formats a custom prompt embedding the question, the student's text, and the official IB markscheme criteria."),
        ("Score & Feedback Streaming", "The AI Tutor streams a simulated score, highlights awarded marking points, and provides clear coaching to achieve a Level 7.")
    ]
    for title, desc in features:
        p = tf.add_paragraph()
        p.text = f"🤖 {title}"
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_before = Pt(8)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Calibri'
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = text_color
        p_desc.space_after = Pt(10)

    # ----------------------------------------------------
    # Slide 4: Revision Mindmaps & Diagrams
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_title(slide, "Core Revision Mindmaps & Diagrams")
    
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    materials = [
        ("Topic 1 Core Revision Mindmap", "Visual ASCII map outline directly in the Guide tab showing links between EVS paradigms, system flows, energy laws, and natural capital."),
        ("EVS System Flow Diagram", "Visual model mapping inputs (culture, media) -> evaluation filter -> outputs (environmental values/policies)."),
        ("System Boundary & Capital Flows", "Annotated boundary exchange models comparing open/closed/isolated systems, alongside capital stock to annual yield replenishment flow rates.")
    ]
    for title, desc in materials:
        p = tf.add_paragraph()
        p.text = f"🧠 {title}"
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_before = Pt(8)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Calibri'
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = text_color
        p_desc.space_after = Pt(10)

    # ----------------------------------------------------
    # Slide 5: IA & EE Research Hub
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_title(slide, "Syllabus Integration: IA & EE Research Hub")
    
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    sections = [
        ("Internal Assessment (IA) Hub", "Breaks down the 30-mark criteria and offers detailed experimental ideas (eutrophication kinetics, soil buffering) with brainstorm prompts for the tutor."),
        ("Extended Essay (EE) Hub", "Outlines the 34-mark interdisciplinary EE criteria and guides research questions (e.g. CCS vs afforestation) mapping Group 3 value analysis with Group 4 scientific measurements.")
    ]
    for title, desc in sections:
        p = tf.add_paragraph()
        p.text = f"🔬 {title}"
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_before = Pt(8)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Calibri'
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = text_color
        p_desc.space_after = Pt(10)

    prs.save("biome_presentation.pptx")
    print("biome_presentation.pptx generated successfully!")

if __name__ == '__main__':
    create_deck()
